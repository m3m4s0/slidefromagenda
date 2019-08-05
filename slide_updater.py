from pptx import Presentation
import sys

def move_slide(presentation, old_index, new_index):
    xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


def read_master_layout(slidedeck):
    layout = slidedeck.slide_master.slide_layouts[4] #select the slide from the Masterlayout you want to use to create new slides
    return layout


def read_agenda_slide(slidedeck):
    agendaslide = slidedeck.slides[1]
    slide_content = []
    for shape in agendaslide.shapes:
        if shape.has_table:
            table = shape.table
            iterrows = iter(table.rows)
            next(iterrows)
            for row in iterrows:
                content = []
                for cell in row.cells:
                    content.append(cell.text)
                slide_content.append(content)
    return slide_content


def create_divider_slides_from_agenda(slidedeck, slidename, agendaslide, insertindex):
    template = read_master_layout(slidedeck)
    slide_content = read_agenda_slide(slidedeck)
    for content in slide_content:
        slide = slidedeck.slides.add_slide(template)
        text = content[1]   # select the title text from agenda
        slide.placeholders[0].text = text
        index = slidedeck.slides.index(slide)
        move_slide(slidedeck, index, insertindex)
        insertindex = insertindex + 1

    slidedeck.save(slidename)


if __name__ == "__main__":
    slidename = sys.argv[1]
    agendaslide = int(sys.argv[2])
    insertindex = int(sys.argv[3])
    slidedeck = Presentation(slidename)
    create_divider_slides_from_agenda(slidedeck, slidename, agendaslide, insertindex)